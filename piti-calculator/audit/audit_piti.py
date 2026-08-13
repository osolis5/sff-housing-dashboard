#!/usr/bin/env python3
"""Audit 1: re-extract every dataset value used by PITI-Calculator from the xlsx
(and the IHS presentation's slide-21 chart) by explicit address, compare against
data.json, and cross-foot internal consistency. Audit 2: validate the
amortization formula against an independent iterative simulation."""
import json, math, sys
import openpyxl

XLSX = '/Users/oscarsolis/Desktop/SFF/Chicago Community Areas_2026_SFF_Indicators_07022026.xlsx'
PPTX = '/Users/oscarsolis/Desktop/SFF/IHS_NLawndale_Dashboard_Presentation_05182026.pptx'
DATA = '/Users/oscarsolis/Desktop/SFF/PITI-Calculator/data.json'

wb = openpyxl.load_workbook(XLSX, read_only=False, data_only=True)
D = json.load(open(DATA))

# Audit 0: the inline <script id="piti-data"> block in index.html must be
# byte-equivalent JSON to the canonical audited data.json (no drift).
import re
for page in ('/Users/oscarsolis/Desktop/SFF/PITI-Calculator/index.html',
             '/Users/oscarsolis/Desktop/SFF/baseline-data-tool/index.html',
             '/Users/oscarsolis/Desktop/SFF/baseline-data-tool/piti-calculator/index.html'):
    html = open(page).read()
    m = re.search(r'<script id="piti-data" type="application/json">\n(.*?)\n</script>', html, re.S)
    inline_ok = m is not None and json.loads(m.group(1)) == D
    name = page.split('/Desktop/SFF/')[1]
    print(f'== INLINE DATA BLOCK: {"PASS" if inline_ok else "FAIL"} — {name} '
          f'{"matches" if inline_ok else "DOES NOT match"} data.json ==')
print()

results = []
def check(label, sheet, cell, expected, tol=0.0):
    actual = wb[sheet][cell].value
    ok = (actual is not None and expected is not None and
          (abs(float(actual) - float(expected)) <= tol if tol else float(actual) == float(expected)))
    results.append((ok, label, f'{sheet}!{cell}', actual, expected))

# ---- G. Mortgage_Lending: price tiers (typical = 2024 median financed value;
#      lower/higher = workbook placeholder cut points at 75% / 125%) ----
check('tier typical (2024 HMDA median)', 'G.Mortgage_Lending', 'G19', D['priceTiers2024']['typical'])
check('tier lower (75% placeholder)',    'G.Mortgage_Lending', 'I19', D['priceTiers2024']['lower'])
check('tier higher (125% placeholder)',  'G.Mortgage_Lending', 'J19', D['priceTiers2024']['higher'])

# ---- E. Property_Sales: all-sales 2024 medians (col U = 2024) + counts ----
check('all-sales median SF 2024',    'E.Property_Sales', 'U124', D['salePrices2024']['sf'])
check('all-sales median condo 2024', 'E.Property_Sales', 'U170', D['salePrices2024']['condo'])
check('all-sales median 2-4 2024',   'E.Property_Sales', 'U216', D['salePrices2024']['u24'])
check('sales count SF 2024',    'E.Property_Sales', 'U109', D['salesCounts2024']['sf'])
check('sales count condo 2024', 'E.Property_Sales', 'U155', D['salesCounts2024']['condo'])
check('sales count 2-4 2024',   'E.Property_Sales', 'U201', D['salesCounts2024']['u24'])

# ---- B. Taxpayer_Characteristics: tax bills (col N = TY2024, M = TY2023,
#      B = TY2012; workbook values carry cents, json stores rounded dollars) ----
check('tax bill SF TY2024',  'B.Taxpayer_Characteristics', 'N18', D['taxBillsTY2024']['sf'], 0.51)
check('tax bill 2-4 TY2024', 'B.Taxpayer_Characteristics', 'N33', D['taxBillsTY2024']['u24'], 0.51)
check('tax bill SF TY2023',  'B.Taxpayer_Characteristics', 'M18', D['taxBillsTY2023']['sf'], 0.51)
check('tax bill 2-4 TY2023', 'B.Taxpayer_Characteristics', 'M33', D['taxBillsTY2023']['u24'], 0.51)
check('tax bill SF TY2012',  'B.Taxpayer_Characteristics', 'B18', D['taxBillsTY2012']['sf'], 0.51)
check('tax bill 2-4 TY2012', 'B.Taxpayer_Characteristics', 'B33', D['taxBillsTY2012']['u24'], 0.51)
check('homeowner exemption SF',  'B.Taxpayer_Characteristics', 'F49', D['exemptions2023']['homeownerShareSF'])
check('homeowner exemption 2-4', 'B.Taxpayer_Characteristics', 'H49', D['exemptions2023']['homeownerShareU24'])
check('senior exemption SF',     'B.Taxpayer_Characteristics', 'F65', D['exemptions2023']['seniorShareSF'])
check('senior exemption 2-4',    'B.Taxpayer_Characteristics', 'H65', D['exemptions2023']['seniorShareU24'])

# ---- A. Demographic_Socioeconomic ----
check('NL median HH income 2024',      'A.Demographic_Socioeconomic', 'C83', D['incomes2024']['nlMedian'],  0.51)
check('NL owner median income 2024',   'A.Demographic_Socioeconomic', 'F83', D['incomes2024']['nlOwner'],   0.51)
check('NL renter median income 2024',  'A.Demographic_Socioeconomic', 'I83', D['incomes2024']['nlRenter'],  0.51)
check('Chicago median income 2024',    'A.Demographic_Socioeconomic', 'C88', D['incomes2024']['chicagoMedian'], 0.51)
check('Chicago renter median income 2024', 'A.Demographic_Socioeconomic', 'I88', D['incomes2024']['chicagoRenter'], 0.51)
check('NL total households 2024', 'A.Demographic_Socioeconomic', 'C67', D['households2024']['total'])
br = D['households2024']['incomeBrackets']
check('bracket <25K share',    'A.Demographic_Socioeconomic', 'F67', br[0]['share'], 1e-4)
check('bracket 25-50K share',  'A.Demographic_Socioeconomic', 'I67', br[1]['share'], 1e-6)
check('bracket 50-100K share', 'A.Demographic_Socioeconomic', 'L67', br[2]['share'], 1e-6)
check('bracket 100K+ share',   'A.Demographic_Socioeconomic', 'O67', br[3]['share'], 1e-6)
check('NL owner share 2024',  'A.Demographic_Socioeconomic', 'I116', D['households2024']['ownerShare'], 1e-6)
check('NL renter share 2024', 'A.Demographic_Socioeconomic', 'O116', D['households2024']['renterShare'], 1e-6)

# ---- G. Mortgage_Lending: loans, values, borrower mixes ----
for cell, key in [('B19','loans2019'), ('C19','medianValue2019'), ('D19','loans2023'),
                  ('E19','medianValue2023'), ('F19','loans2024'), ('G19','medianValue2024')]:
    check(f'HMDA {key}', 'G.Mortgage_Lending', cell, D['hmda'][key])
inc = D['hmda']['incomeMix2324']
for i, cell in enumerate(['F67', 'G67', 'H67', 'I67']):
    check(f'HMDA income count {inc[i]["label"]}', 'G.Mortgage_Lending', cell, inc[i]['count'])
for i, cell in enumerate(['F83', 'G83', 'H83', 'I83']):
    check(f'HMDA income share {inc[i]["label"]}', 'G.Mortgage_Lending', cell, inc[i]['share'], 1e-6)
check('HMDA upper-income share 2019', 'G.Mortgage_Lending', 'E83', D['hmda']['incomeMix2019Upper'], 1e-6)
race = D['hmda']['raceMix2324']
rcells = [('I35','I51'), ('J35','J51'), ('H35','H51'), ('K35','K51'), ('M35','M51')]  # AA, Hisp, White, Asian, NA
for i, (ccell, scell) in enumerate(rcells):
    check(f'HMDA race count {race[i]["label"]}', 'G.Mortgage_Lending', ccell, race[i]['count'])
    check(f'HMDA race share {race[i]["label"]}', 'G.Mortgage_Lending', scell, race[i]['share'], 1e-6)

# ---- IHS presentation slide 21: median homebuyer incomes (HMDA) ----
# Not in the workbook — read the values straight from the slide's chart data.
pres_results = []
try:
    from pptx import Presentation
    slide = Presentation(PPTX).slides[20]              # slide 21, 0-indexed
    chart = next(s.chart for s in slide.shapes if s.has_chart)
    cats = [str(c) for c in chart.plots[0].categories]
    series = {s.name: dict(zip(cats, s.values)) for s in chart.plots[0].series}
    for label, cat, year, key in [
        ('NL homebuyer income 2024',      'North Lawndale',  '2024', 'nl'),
        ('Chicago homebuyer income 2024', 'City of Chicago', '2024', 'chicago'),
        ('NL homebuyer income 2019',      'North Lawndale',  '2019', 'nl2019'),
        ('Chicago homebuyer income 2019', 'City of Chicago', '2019', 'chicago2019'),
    ]:
        actual = series[year][cat]
        expected = D['homebuyerIncomes2024'][key]
        pres_results.append((float(actual) == float(expected), label, f'slide21 {year}/{cat}', actual, expected))
except Exception as exc:  # missing pptx / python-pptx — report, do not crash
    pres_results.append((False, 'slide-21 chart read', PPTX, f'ERROR: {exc}', 'readable chart'))

# ---- cross-foot internal consistency ----
xf = []
xf.append(('tier lower = 75% of typical (placeholder rule)',
           D['priceTiers2024']['lower'], 0.75 * D['priceTiers2024']['typical']))
xf.append(('tier higher = 125% of typical (placeholder rule)',
           D['priceTiers2024']['higher'], 1.25 * D['priceTiers2024']['typical']))
xf.append(('typical tier = HMDA 2024 median value',
           D['priceTiers2024']['typical'], D['hmda']['medianValue2024']))
xf.append(('income-mix counts sum = 2023+2024 loans',
           sum(r['count'] for r in inc), D['hmda']['loans2023'] + D['hmda']['loans2024']))
xf.append(('race-mix counts sum = 2023+2024 loans',
           sum(r['count'] for r in race), D['hmda']['loans2023'] + D['hmda']['loans2024']))
xf.append(('income bracket shares sum to 1', round(sum(b['share'] for b in br), 6), 1.0))
xf.append(('owner+renter share = 1', round(D['households2024']['ownerShare'] + D['households2024']['renterShare'], 6), 1.0))
own = wb['A.Demographic_Socioeconomic']['F116'].value
rent = wb['A.Demographic_Socioeconomic']['L116'].value
xf.append(('owner+renter households = total', own + rent, D['households2024']['total']))
xf.append(('2019 HMDA income counts sum = 2019 loans',
           sum(wb['G.Mortgage_Lending'][c].value for c in ['B67','C67','D67','E67']), D['hmda']['loans2019']))
q18 = wb['B.Taxpayer_Characteristics']['Q18'].value   # TY2023→TY2024 pct change, SF
xf.append(('SF tax-bill change TY2023→TY2024 ≈ +73% (About-card claim)',
           round(q18, 2), 0.73))

# ---- Audit 2: amortization formula vs independent iterative simulation ----
def pni_closed(loan, apr, years):
    r, n = apr / 100 / 12, years * 12
    return loan / n if r == 0 else loan * r * (1 + r) ** n / ((1 + r) ** n - 1)

def sim_balance(loan, apr, years, payment):
    """simulate month-by-month; return ending balance (should be ~0)"""
    r, bal = apr / 100 / 12, loan
    for _ in range(years * 12):
        bal = bal * (1 + r) - payment
    return bal

fchecks = []
for loan, apr, yrs in [(100000, 6.0, 30), (200000, 6.0, 30), (300000, 7.0, 30),
                       (308750, 6.5, 30), (231562.5, 6.5, 30), (385937.5, 6.5, 30),
                       (308750, 6.5, 15), (95000, 3.0, 20), (760000, 10.0, 30)]:
    m = pni_closed(loan, apr, yrs)
    endbal = sim_balance(loan, apr, yrs, m)
    fchecks.append((loan, apr, yrs, m, endbal, abs(endbal) < 0.01))
# textbook reference values
refs = [(100000, 6.0, 30, 599.55), (200000, 6.0, 30, 1199.10), (300000, 7.0, 30, 1995.91)]
refchecks = [(l, a, y, ref, pni_closed(l, a, y), abs(pni_closed(l, a, y) - ref) < 0.01) for l, a, y, ref in refs]

# ---- report ----
results_all = results + pres_results
fails = [r for r in results_all if not r[0]]
print(f'== DATASET AUDIT: {len(results_all)} value checks, {len(results_all)-len(fails)} PASS, {len(fails)} FAIL ==')
for ok, label, ref, actual, expected in results:
    print(f'  {"PASS" if ok else "FAIL"}  {label:42s} {ref:32s} xlsx={actual}  json={expected}')
for ok, label, ref, actual, expected in pres_results:
    print(f'  {"PASS" if ok else "FAIL"}  {label:42s} {ref:32s} pptx={actual}  json={expected}')
print(f'\n== CROSS-FOOT CONSISTENCY ==')
xf_fails = 0
for label, a, b in xf:
    ok = abs(a - b) < 1e-9
    xf_fails += 0 if ok else 1
    print(f'  {"PASS" if ok else "FAIL"}  {label:52s} {a} vs {b}')
print(f'\n== FORMULA: closed-form vs 360-payment simulation (ending balance ~ $0) ==')
for loan, apr, yrs, m, endbal, ok in fchecks:
    print(f'  {"PASS" if ok else "FAIL"}  ${loan:>9,.0f} @ {apr}% {yrs}yr -> M=${m:,.2f}, end balance ${endbal:+.6f}')
print(f'\n== FORMULA: vs published reference payments ==')
for l, a, y, ref, got, ok in refchecks:
    print(f'  {"PASS" if ok else "FAIL"}  ${l:,} @ {a}% {y}yr: reference ${ref} vs computed ${got:.2f}')
sys.exit(1 if (fails or xf_fails) else 0)
